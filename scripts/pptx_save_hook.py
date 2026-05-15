"""Hook script: runs after mcp__ppt__save_presentation to remove empty placeholders."""
import sys
import json

def main():
    data = json.loads(sys.stdin.buffer.read().decode("utf-8-sig"))
    file_path = data.get("tool_input", {}).get("file_path", "")
    if not file_path or not file_path.endswith(".pptx"):
        return

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    removed = 0
    for slide in prs.slides:
        to_remove = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
                continue
            if not shape.has_text_frame or shape.text_frame.text.strip() == "":
                to_remove.append(shape)
        for shape in to_remove:
            shape._element.getparent().remove(shape._element)
            removed += 1
    if removed:
        prs.save(file_path)
        print(f"Removed {removed} empty placeholder(s) from {file_path}", flush=True)

if __name__ == "__main__":
    main()
