"""Remove empty placeholders that show 'クリックしてテキストを追加' from a PPTX file."""
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def remove_empty_placeholders(path: str) -> None:
    prs = Presentation(path)

    for slide_idx, slide in enumerate(prs.slides):
        to_remove = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
                continue
            # Remove placeholder only if it contains no text
            if not shape.has_text_frame:
                to_remove.append(shape)
                continue
            text = shape.text_frame.text.strip()
            if text == "":
                to_remove.append(shape)

        for shape in to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"  Slide {slide_idx}: removed placeholder '{shape.name}'")

    prs.save(path)
    print(f"Saved: {path}")


if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else r"output\sample_四季.pptx"
    remove_empty_placeholders(target)
